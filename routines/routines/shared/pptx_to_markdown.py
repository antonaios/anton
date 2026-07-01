"""PPTX → Markdown converter — a standalone port of microsoft/markitdown's
``_pptx_converter.py`` (NO markitdown runtime dependency).

Lets ANTON ingest a PowerPoint deck (e.g. a competitor's pitch dropped into a
workspace's VDR) as token-efficient markdown for the recall + chat lanes, and
backs the future #27 ``/pitch`` composite.

Design (preserved from markitdown, with ANTON adaptations):
  * Shapes are iterated in **position-sorted order** ((top, left), top-to-bottom
    then left-to-right) so the markdown reading order matches the slide's visual
    order — naive python-pptx iteration returns shapes in z-order, which
    scrambles multi-column slides.
  * Slide titles → ``## Slide N: <title>`` headings (markitdown used ``#``; we
    number the slides for recall-friendly anchoring).
  * Body text frames → markdown paragraphs (single paragraph) or bullet lists
    (multiple paragraphs, indented by paragraph level).
  * Tables → markdown tables (built directly from the cell grid; markitdown
    round-trips through HTML + ``markdownify`` — we avoid that dependency).
  * Charts → a markdown table of the chart's underlying data
    (``chart.plots[0].categories`` × ``chart.series`` values); unsupported
    chart types degrade to ``[unsupported chart]``.
  * Grouped shapes → recursive descent (also position-sorted).
  * Speaker notes → a ``### Notes:`` block under each slide.
  * Pictures → a ``![<alt or "image">]`` placeholder. We deliberately drop
    markitdown's ``data:image/png;base64,…`` inlining (far too token-heavy).
    Optional image *description* is OFF by default; when enabled it routes
    through ANTON's sensitivity-aware LLM router instead of markitdown's raw
    ``llm_client`` — and a confidential/MNPI deck is pinned to the LOCAL
    (Ollama) lane, never cloud (CLAUDE.md #no-mnpi-to-cloud — was cited as §5.4).

The default path needs NO LLM — it is pure-CPU python-pptx text/table/chart
extraction.

----------------------------------------------------------------------------
Attribution: logic ported from microsoft/markitdown ``_pptx_converter.py``
(https://github.com/microsoft/markitdown), MIT License,
Copyright (c) Microsoft Corporation. We lift the converter logic only; we do
NOT take markitdown as a runtime dependency (its other converters are either
redundant with or worse than ANTON's existing paths — see
``evaluations/MARKITDOWN-EVALUATION-2026-05-28.md``).
----------------------------------------------------------------------------
"""

from __future__ import annotations

import base64
import ipaddress
import logging
import os
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import IO, Any
from urllib.parse import urlparse

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

from routines.shared.ollama_client import DEFAULT_BASE_URL, OllamaClient
from routines.shared.routing import Sensitivity, lane_to_model, pick_lane

log = logging.getLogger(__name__)


class PptxConversionError(RuntimeError):
    """Raised when a .pptx cannot be converted (bad file, or a sensitivity
    invariant would be violated)."""


# ── preflight caps (zip-bomb / oversized-file guard) ──────────────────────────
# A legit inbound pitch deck is a few MB to a few tens of MB on disk and inflates
# to well under these ceilings; an Office-XML archive compresses ≈4–15×. These
# caps stop a decompression bomb (a tiny archive that inflates to GBs, or one
# with a pathological member count) from being eagerly inflated by python-pptx
# into memory. Tuned generously so no real deck trips them.
_MAX_ARCHIVE_BYTES = 200 * 1024 * 1024          # 200 MB compressed-on-disk
_MAX_UNCOMPRESSED_BYTES = 1024 * 1024 * 1024    # 1 GB inflated total
_MAX_MEMBERS = 10_000                           # zip entry count
# CPython's zipfile walks the central directory by its (Zip64-resolved) SIZE, not
# the declared entry count, building a ZipInfo per header — so the directory is
# the real lever for an entry-count amplification attack. We bound it two ways:
# a cheap size cap here (also the backstop for the exotic prepended-data case the
# header counter can't follow), plus a tight bounded header count below. 4 MB
# comfortably holds a legitimate ≤_MAX_MEMBERS-entry directory (a real entry is
# ~80–160 B, so 10k entries ≈ 1.6 MB) while bounding a hostile one.
_MAX_CENTRAL_DIR_BYTES = 4 * 1024 * 1024
_MAX_COMPRESSION_RATIO = 200                    # inflated / compressed (aggregate AND per-member)
# Per-member ratio is only enforced for members that inflate past this floor: a
# small but highly-repetitive XML part can legitimately compress >200×, and a
# small member is harmless however it compresses. A bomb member inflates to
# hundreds of MB, far past the floor.
_MEMBER_RATIO_FLOOR = 4 * 1024 * 1024           # 4 MB inflated


# One concise caption per image; kept short to stay token-light. Only used when
# image description is explicitly enabled.
_IMAGE_DESCRIBE_PROMPT = (
    "Describe this slide image in one concise sentence for a reading-order "
    "transcript. State what it depicts; do not speculate beyond what is shown."
)


@dataclass(frozen=True)
class _Opts:
    describe_images: bool
    sensitivity: Sensitivity
    ollama_client: OllamaClient | None


# ── public entry point ──────────────────────────────────────────────────────


def pptx_to_markdown(
    source: str | Path | IO[bytes],
    *,
    describe_images: bool = False,
    sensitivity: Sensitivity = "confidential",
    ollama_base_url: str = DEFAULT_BASE_URL,
) -> str:
    """Convert a ``.pptx`` to markdown.

    Args:
        source: path (``str``/``Path``) or an open binary file-like object.
        describe_images: when True, pictures are captioned via the LOCAL
            multimodal lane; when False (default) pictures emit a text
            placeholder and NO LLM is touched.
        sensitivity: routing sensitivity for image description. Every tier is
            pinned to the LOCAL lane — image bytes never leave the box. Defaults
            to ``confidential`` (the safe assumption for an inbound deck of
            unknown provenance).
        ollama_base_url: base URL for the local Ollama used for image
            description. It MUST be loopback for EVERY sensitivity — a slide
            image is never sent to a non-loopback host (CLAUDE.md
            #no-mnpi-to-cloud — was cited as §5.4).

    Returns:
        The deck as a single markdown string (one ``## Slide N`` section each).
    """
    client: OllamaClient | None = None
    if describe_images:
        # Hard rule (CLAUDE.md #no-mnpi-to-cloud — was cited as §5.4): a slide
        # image's bytes must not leave the box, regardless of sensitivity. We
        # do NOT accept a caller-supplied
        # client object — a subclass could lie about its endpoint and override
        # chat() to exfiltrate. Instead we construct a concrete OllamaClient
        # ourselves and require the endpoint to be loopback for ALL tiers (the
        # lane guard in _describe_image proves the *provider* is ollama; this
        # proves the *destination* is on-box). There is no operator opt-out to a
        # remote Ollama — image description is local-only or off.
        if not _is_loopback_url(ollama_base_url):
            raise PptxConversionError(
                "image description requires a loopback Ollama endpoint; got "
                f"{ollama_base_url!r} — refusing to send a slide image off-box"
            )
        # A loopback URL is necessary but not sufficient: OllamaClient calls
        # requests, which honours HTTP_PROXY/ALL_PROXY for a 127.0.0.1 request
        # unless NO_PROXY exempts it — i.e. the image bytes could be tunnelled to
        # a remote proxy despite the loopback host. Fail CLOSED if the environment
        # would route this URL through any proxy (CLAUDE.md #no-mnpi-to-cloud
        # — was cited as §5.4; absolute).
        if _proxy_would_intercept(ollama_base_url):
            raise PptxConversionError(
                "a proxy is configured for the Ollama endpoint "
                f"({ollama_base_url!r}); refusing image description so slide image "
                "bytes stay on-box — unset HTTP(S)_PROXY/ALL_PROXY or add the host "
                "to NO_PROXY"
            )
        client = OllamaClient(base_url=ollama_base_url)

    opts = _Opts(
        describe_images=describe_images,
        sensitivity=sensitivity,
        ollama_client=client,
    )

    # For a path, open ONE handle and use it for BOTH the preflight and the parse,
    # so the file can't be swapped for a bomb between inspection and parse
    # (TOCTOU); close it only after rendering completes. A file-like source is
    # inspected + parsed in place.
    if isinstance(source, (str, Path)):
        try:
            handle = open(source, "rb")  # noqa: SIM115 — closed in the finally below
        except OSError as e:
            raise PptxConversionError(f"could not open .pptx: {e}") from e
        try:
            return _convert_stream(handle, opts)
        finally:
            handle.close()
    return _convert_stream(source, opts)


def _convert_stream(stream: IO[bytes], opts: _Opts) -> str:
    """Preflight, open, and render a seekable binary .pptx stream. The caller owns
    the stream's lifetime (we never close it)."""
    # Reject a hostile/oversized archive BEFORE python-pptx eagerly inflates it.
    _preflight_zip(stream)
    try:
        prs = Presentation(stream)
    except PptxConversionError:
        raise
    except Exception as e:  # noqa: BLE001 — any pptx/zip/parse failure → clean error
        raise PptxConversionError(f"could not open .pptx: {e}") from e

    # One bad slide degrades to a placeholder rather than aborting the deck; a
    # PptxConversionError (the sensitivity invariant) is never swallowed.
    try:
        n_slides = len(prs.slides)
    except Exception as e:  # noqa: BLE001
        raise PptxConversionError(f"could not read slides: {e}") from e
    slides_md: list[str] = []
    for idx in range(1, n_slides + 1):
        try:
            slides_md.append(_convert_slide(idx, prs.slides[idx - 1], opts))
        except PptxConversionError:
            raise
        except Exception as e:  # noqa: BLE001 — degrade this slide, keep the deck
            log.warning("pptx: slide %d failed to render (%s)", idx, e)
            slides_md.append(f"## Slide {idx}\n\n[slide could not be rendered]")
    body = "\n\n".join(s for s in slides_md if s).strip()
    return body + "\n" if body else ""


# ── preflight ─────────────────────────────────────────────────────────────────


def _preflight_zip(stream: IO[bytes]) -> None:
    """Inspect a binary .pptx *stream*'s zip central directory and reject it
    before handing it to python-pptx if it breaches a size/shape cap. No member
    is decompressed — only ``ZipInfo`` metadata (compressed/uncompressed sizes,
    count) is read, so this cannot itself be bombed. The total archive size is
    checked FIRST, before the central directory is even parsed, so a grossly
    oversized archive is rejected without allocating per-member metadata. Raises
    :class:`PptxConversionError` on breach, a bad/unreadable zip, or a
    non-seekable stream; leaves the stream rewound to where it began."""
    # Absolute archive size, measured from the start of the stream (which is what
    # zipfile and python-pptx read — the central directory lives at the end), then
    # rewind. A non-seekable stream is rejected up front: zipfile and python-pptx
    # both require seeking, so continuing would surface a raw exception from deep
    # inside instead of a clean PptxConversionError.
    try:
        start_pos = stream.tell()
        stream.seek(0, os.SEEK_END)
        archive_bytes = stream.tell()
        stream.seek(start_pos)
    except (OSError, AttributeError) as e:
        raise PptxConversionError(
            "file-like .pptx input must be seekable for preflight inspection"
        ) from e

    if archive_bytes > _MAX_ARCHIVE_BYTES:
        raise PptxConversionError(
            f"rejected .pptx: archive {archive_bytes} bytes exceeds cap {_MAX_ARCHIVE_BYTES}"
        )

    # Bound the work zipfile.ZipFile() will do BEFORE calling it — it eagerly
    # walks the central directory (by its Zip64-resolved SIZE, not the declared
    # entry count) and materializes a ZipInfo per header, so a hostile directory
    # with millions of entries would OOM us inside the parse itself. From the EOCD
    # (Zip64-aware) we get the count, size, and offset and: (1) reject an
    # over-cap declared count or directory size cheaply; (2) tightly count the
    # actual headers up to _MAX_MEMBERS+1 (zipfile ignores the declared count, so
    # a small-count/large-directory archive must be caught by the real count).
    info = _eocd_central_dir_info(stream, archive_bytes)
    try:
        stream.seek(start_pos)
    except (OSError, AttributeError):
        pass
    if info is not None:
        declared_count, cd_size, cd_offset = info
        if declared_count > _MAX_MEMBERS:
            raise PptxConversionError(
                f"rejected .pptx: {declared_count} declared zip members exceeds cap {_MAX_MEMBERS}"
            )
        if cd_size > _MAX_CENTRAL_DIR_BYTES:
            raise PptxConversionError(
                f"rejected .pptx: central directory {cd_size} bytes exceeds cap "
                f"{_MAX_CENTRAL_DIR_BYTES} (too many members?)"
            )
        actual = _count_central_dir_headers(stream, cd_offset, cd_size, _MAX_MEMBERS)
        try:
            stream.seek(start_pos)
        except (OSError, AttributeError):
            pass
        if actual is not None and actual > _MAX_MEMBERS:
            raise PptxConversionError(
                f"rejected .pptx: central directory holds more than {_MAX_MEMBERS} "
                "members (decompression bomb?)"
            )

    try:
        zf = zipfile.ZipFile(stream)
    except zipfile.BadZipFile as e:
        raise PptxConversionError(f"not a valid .pptx (zip) archive: {e}") from e
    except Exception as e:  # noqa: BLE001 — any unreadable/odd source → clean error
        raise PptxConversionError(f"could not read .pptx: {e}") from e
    try:
        infos = zf.infolist()
        member_count = len(infos)
        total_compressed = sum(i.compress_size for i in infos)
        total_uncompressed = sum(i.file_size for i in infos)
        # Per-member ratio: a bomb can pad the archive with incompressible junk to
        # dilute the AGGREGATE ratio while one referenced member still expands
        # pathologically. Flag any large-enough member whose own ratio is extreme.
        worst_member: tuple[str, float] | None = None
        for i in infos:
            if i.file_size > _MEMBER_RATIO_FLOOR:
                member_ratio = i.file_size / max(i.compress_size, 1)
                if member_ratio > _MAX_COMPRESSION_RATIO and (
                    worst_member is None or member_ratio > worst_member[1]
                ):
                    worst_member = (i.filename, member_ratio)
    finally:
        # A ZipFile opened from a passed-in file object leaves the object open;
        # we rewind it below for the subsequent Presentation() read.
        zf.close()

    try:
        stream.seek(start_pos)
    except (OSError, AttributeError):
        pass

    if member_count > _MAX_MEMBERS:
        raise PptxConversionError(
            f"rejected .pptx: {member_count} zip members exceeds cap {_MAX_MEMBERS}"
        )
    if total_uncompressed > _MAX_UNCOMPRESSED_BYTES:
        raise PptxConversionError(
            f"rejected .pptx: uncompressed {total_uncompressed} bytes exceeds cap "
            f"{_MAX_UNCOMPRESSED_BYTES} (decompression bomb?)"
        )
    if worst_member is not None:
        raise PptxConversionError(
            f"rejected .pptx: member {worst_member[0]!r} compression ratio "
            f"{worst_member[1]:.0f}× exceeds cap {_MAX_COMPRESSION_RATIO}× "
            "(decompression bomb?)"
        )
    ratio = total_uncompressed / max(total_compressed, 1)
    if ratio > _MAX_COMPRESSION_RATIO:
        raise PptxConversionError(
            f"rejected .pptx: aggregate compression ratio {ratio:.0f}× exceeds cap "
            f"{_MAX_COMPRESSION_RATIO}× (decompression bomb?)"
        )


# Zip End-Of-Central-Directory record (PK\x05\x06): 16-bit total entries at
# offset 10, 32-bit central-directory size at offset 12, 32-bit central-directory
# offset at offset 16. A Zip64 EOCD locator (PK\x06\x07) sitting 20 bytes before
# the EOCD points to a Zip64 EOCD (PK\x06\x06) whose 64-bit total/size/offset are
# at offsets 32/40/48. CPython's zipfile consults the Zip64 record whenever the
# locator is present (NOT only when the legacy fields are the 0xFFFF/0xFFFFFFFF
# sentinels) and lets it OVERRIDE the legacy values — so we must do the same or a
# small-legacy/large-Zip64 archive would slip past the caps.
_EOCD_SIG = b"PK\x05\x06"
_ZIP64_LOCATOR_SIG = b"PK\x06\x07"
_ZIP64_EOCD_SIG = b"PK\x06\x06"
_MAX_ZIP_COMMENT = 0xFFFF
_CDH_SIG = b"PK\x01\x02"   # central-directory file header
_CDH_FIXED = 46            # fixed-size portion of a central-directory header


def _eocd_central_dir_info(stream: IO[bytes], archive_bytes: int) -> tuple[int, int, int] | None:
    """Read ``(entry_count, central_directory_size, central_directory_offset)``
    from the zip's EOCD — the same values CPython's zipfile will act on — WITHOUT
    parsing the central directory, so a pathological directory is caught before
    zipfile walks it and allocates a ``ZipInfo`` per header. Returns ``None``
    (caller falls back to the post-parse cap) when the EOCD can't be located.

    Mirrors CPython: if a Zip64 locator is present its record is AUTHORITATIVE
    (overrides the legacy fields). If a locator is present but its record can't be
    read, or a legacy field is a sentinel with no resolving record, we fail CLOSED
    — report cap-tripping values so the caller rejects rather than letting zipfile
    walk an unbounded directory."""
    tail_len = min(archive_bytes, 22 + _MAX_ZIP_COMMENT)
    try:
        stream.seek(archive_bytes - tail_len)
        tail = stream.read(tail_len)
    except (OSError, ValueError):
        return None

    pos = _find_eocd(tail)
    if pos is None:
        return None
    count = int.from_bytes(tail[pos + 10 : pos + 12], "little")
    cd_size = int.from_bytes(tail[pos + 12 : pos + 16], "little")
    cd_offset = int.from_bytes(tail[pos + 16 : pos + 20], "little")

    has_locator = pos >= 20 and tail[pos - 20 : pos - 16] == _ZIP64_LOCATOR_SIG
    if has_locator:
        # The locator makes the Zip64 record authoritative for zipfile. Use it, or
        # fail closed if it can't be read (zipfile itself would raise BadZipFile).
        z64 = _read_zip64_eocd(stream, tail, pos)
        if z64 is not None:
            return z64
        return (_MAX_MEMBERS + 1, _MAX_CENTRAL_DIR_BYTES + 1, cd_offset)
    if count == 0xFFFF or cd_size == 0xFFFFFFFF:
        # A sentinel field with no locator is malformed → fail closed on it.
        return (
            _MAX_MEMBERS + 1 if count == 0xFFFF else count,
            _MAX_CENTRAL_DIR_BYTES + 1 if cd_size == 0xFFFFFFFF else cd_size,
            cd_offset,
        )
    return (count, cd_size, cd_offset)


def _read_zip64_eocd(stream: IO[bytes], tail: bytes, eocd_pos: int) -> tuple[int, int, int] | None:
    """``(total_entries, central_directory_size, central_directory_offset)`` from
    the Zip64 EOCD record that the locator (20 bytes before the EOCD) points to, or
    ``None`` if the locator/record can't be read."""
    loc = eocd_pos - 20
    if loc < 0 or tail[loc : loc + 4] != _ZIP64_LOCATOR_SIG:
        return None
    z64_ofs = int.from_bytes(tail[loc + 8 : loc + 16], "little")
    try:
        stream.seek(z64_ofs)
        rec = stream.read(56)
    except (OSError, ValueError):
        return None
    if len(rec) < 56 or rec[:4] != _ZIP64_EOCD_SIG:
        return None
    return (
        int.from_bytes(rec[32:40], "little"),
        int.from_bytes(rec[40:48], "little"),
        int.from_bytes(rec[48:56], "little"),
    )


def _count_central_dir_headers(
    stream: IO[bytes], cd_offset: int, cd_size: int, limit: int
) -> int | None:
    """Count central-directory file headers, bounded to ``limit + 1``, by walking
    the fixed 46-byte headers only (the variable name/extra/comment are skipped via
    seek) — NO ``ZipInfo`` is built. This bounds the real entry count before
    zipfile (which walks the directory by size, ignoring the declared count) gets
    a chance to allocate millions of entries. Returns the count (capped at
    ``limit + 1``), or ``None`` if the directory can't be read."""
    try:
        stream.seek(cd_offset)
    except (OSError, ValueError):
        return None
    count = 0
    consumed = 0
    while consumed < cd_size and count <= limit:
        hdr = stream.read(_CDH_FIXED)
        if len(hdr) < _CDH_FIXED or hdr[:4] != _CDH_SIG:
            break  # end of the central directory (or malformed) — stop counting
        n = int.from_bytes(hdr[28:30], "little")
        m = int.from_bytes(hdr[30:32], "little")
        k = int.from_bytes(hdr[32:34], "little")
        try:
            stream.seek(n + m + k, os.SEEK_CUR)
        except (OSError, ValueError):
            return None
        consumed += _CDH_FIXED + n + m + k
        count += 1
    return count


def _find_eocd(tail: bytes) -> int | None:
    """Index of the EOCD record CPython's ``zipfile`` will act on. We deliberately
    mirror ``zipfile._EndRecData``: the LAST ``PK\\x05\\x06`` signature in the tail
    for which the 22-byte fixed record fits. We do NOT additionally require the
    declared comment length to consume the rest of the file — ``zipfile`` doesn't,
    so demanding it would make the preflight MISS an EOCD that ``zipfile`` still
    honours (e.g. with trailing bytes appended), silently skipping every cap while
    ``zipfile`` goes on to parse the archive. Reading the same record ``zipfile``
    reads keeps the guard and the parser consistent."""
    pos = tail.rfind(_EOCD_SIG)
    if pos < 0 or pos + 22 > len(tail):
        return None
    return pos


# ── slide / shape rendering ───────────────────────────────────────────────────


def _convert_slide(idx: int, slide: Any, opts: _Opts) -> str:
    title_shape = _slide_title(slide)
    title_text = ""
    if title_shape is not None:
        try:
            if title_shape.has_text_frame:
                title_text = _escape_inline(title_shape.text_frame.text)
        except Exception as e:  # noqa: BLE001 — a bad title degrades to no title
            log.debug("pptx: slide %d title unreadable (%s)", idx, e)

    blocks: list[str] = [f"## Slide {idx}: {title_text}" if title_text else f"## Slide {idx}"]

    try:
        shapes = _sorted_shapes(slide.shapes)
    except Exception as e:  # noqa: BLE001 — shapes unreadable → just the heading
        log.warning("pptx: slide %d shapes unreadable (%s)", idx, e)
        shapes = []
    for shape in shapes:
        try:
            blocks.extend(_render_shape(shape, title_shape, opts))
        except PptxConversionError:
            raise  # sensitivity invariant — never degrade past this
        except Exception as e:  # noqa: BLE001 — one bad shape degrades, slide survives
            log.debug("pptx: shape skipped (%s)", e)
            blocks.append("[unreadable shape]")

    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text.strip() if notes_tf is not None else ""
            if notes:
                blocks.append("### Notes:\n\n" + notes)
    except Exception as e:  # noqa: BLE001 — notes are best-effort
        log.debug("pptx: slide %d notes skipped (%s)", idx, e)

    return "\n\n".join(b for b in blocks if b).strip()


def _render_shape(shape: Any, title_shape: Any, opts: _Opts) -> list[str]:
    """Render a single shape to zero-or-more markdown blocks (recursing into
    groups). The slide's title shape is skipped — it is already the heading.
    Compared by XML element, not object identity: python-pptx hands back a
    fresh wrapper per access, so ``shape is title_shape`` is unreliable."""
    if title_shape is not None and shape._element is title_shape._element:
        return []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        out: list[str] = []
        for sub in _sorted_shapes(shape.shapes):
            try:
                out.extend(_render_shape(sub, title_shape, opts))
            except PptxConversionError:
                raise  # sensitivity invariant — never degrade past this
            except Exception as e:  # noqa: BLE001 — one bad grouped shape degrades
                log.debug("pptx: grouped shape skipped (%s)", e)
                out.append("[unreadable shape]")
        return out

    if _is_picture(shape):
        return [_picture_md(shape, opts)]

    if getattr(shape, "has_table", False):
        md = _table_md(shape.table)
        return [md] if md else []

    if getattr(shape, "has_chart", False):
        return [_chart_md(shape.chart)]

    if shape.has_text_frame:
        md = _text_frame_md(shape.text_frame, is_body_placeholder=_is_body_placeholder(shape))
        return [md] if md else []

    return []


# Placeholder types whose text frames are conventionally bulleted outlines. A
# title/subtitle or a plain (non-placeholder) text box is prose by default.
_BODY_PLACEHOLDER_TYPES = frozenset(
    {
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.VERTICAL_BODY,
        PP_PLACEHOLDER.VERTICAL_OBJECT,
    }
)


def _is_body_placeholder(shape: Any) -> bool:
    """True only for a body/content placeholder — the case where PowerPoint
    inherits bullet formatting from the layout/master that python-pptx can't
    resolve. Used to gate the implicit-bullet heuristic so a plain text box's
    multiple paragraphs are NOT auto-bulleted."""
    try:
        if not shape.is_placeholder:
            return False
        return shape.placeholder_format.type in _BODY_PLACEHOLDER_TYPES
    except (AttributeError, ValueError, KeyError):
        return False


def _sorted_shapes(shapes: Any) -> list[Any]:
    """Position-sort shapes top-to-bottom then left-to-right. ``top``/``left``
    can be ``None`` for some shapes (no explicit offset) — treat as 0."""
    return sorted(
        shapes,
        key=lambda s: (s.top if s.top is not None else 0, s.left if s.left is not None else 0),
    )


def _slide_title(slide: Any) -> Any | None:
    try:
        return slide.shapes.title
    except (AttributeError, KeyError, ValueError):
        return None


# ── text ──────────────────────────────────────────────────────────────────────


# Markdown inline specials that would break a heading line or the ``![...]``
# image-alt syntax these values are dropped into.
_INLINE_SPECIALS = re.compile(r"([\\`*_\[\]])")


def _escape_inline(text: str) -> str:
    """Normalize a title / alt-text for safe inline use: collapse all whitespace
    (incl. newlines, which would otherwise split a heading line) to single
    spaces, strip, then backslash-escape markdown specials. Untrusted deck text
    can't smuggle a ``]`` to close an image alt early or ``*``/``_`` emphasis
    into a heading."""
    if not text:
        return ""
    collapsed = " ".join(text.split())
    return _INLINE_SPECIALS.sub(r"\\\1", collapsed)


def _text_frame_md(tf: Any, *, is_body_placeholder: bool) -> str:
    """Render a text frame as paragraphs and/or a bullet list.

    Bullet state is read from each paragraph's XML when it is stated explicitly
    (``a:buChar``/``a:buAutoNum`` → bullet; ``a:buNone`` → plain). PowerPoint
    most often *inherits* bullet formatting from the layout/master, which
    python-pptx does not resolve. We only fall back to the implicit-bullet
    heuristic — an indented paragraph, or any paragraph in a multi-paragraph
    frame, is a bullet — when the shape is a body/content placeholder (where that
    inheritance applies). A plain text box (or a title/subtitle) is treated as
    prose: its paragraphs render blank-line-separated unless they carry explicit
    bullet XML, so a multi-line text box isn't mis-rendered as a list just
    because it has >1 paragraph. markitdown emitted the raw ``shape.text`` with
    no bullet structure at all."""
    paras = [(p.text.strip(), p.level or 0, _para_bullet(p)) for p in tf.paragraphs]
    paras = [(t, lvl, b) for t, lvl, b in paras if t]
    if not paras:
        return ""

    has_explicit = any(b is not None for _, _, b in paras)
    multi = len(paras) > 1

    entries: list[tuple[bool, str]] = []  # (is_bullet, rendered_line)
    for text, lvl, bullet in paras:
        if bullet is not None:
            is_bullet = bullet
        elif is_body_placeholder:
            is_bullet = multi or lvl > 0
        else:
            is_bullet = False
        if is_bullet:
            entries.append((True, f"{'  ' * max(lvl, 0)}- {text}"))
        else:
            entries.append((False, text))

    # A single, plainly-formatted paragraph reads better without a list marker.
    if len(entries) == 1 and not has_explicit and paras[0][1] == 0:
        return paras[0][0]

    # Separate markdown blocks with a blank line so a plain paragraph adjacent to
    # a bullet isn't swallowed as list-item continuation text. Consecutive
    # bullets stay tight (one list).
    out: list[str] = []
    prev_bullet: bool | None = None
    for is_bullet, line in entries:
        if out and not (is_bullet and prev_bullet):
            out.append("")
        out.append(line)
        prev_bullet = is_bullet
    return "\n".join(out)


def _para_bullet(paragraph: Any) -> bool | None:
    """Explicit bullet state of a paragraph from its ``a:pPr``:
    ``True`` if it declares a bullet, ``False`` if it declares no-bullet,
    ``None`` if it states neither (formatting inherited from the layout)."""
    try:
        pPr = paragraph._p.find(qn("a:pPr"))
        if pPr is None:
            return None
        if pPr.find(qn("a:buNone")) is not None:
            return False
        if pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None:
            return True
    except Exception:  # noqa: BLE001 — bullet detection is best-effort
        return None
    return None


# ── tables ──────────────────────────────────────────────────────────────────


def _table_md(table: Any) -> str:
    rows = list(table.rows)
    if not rows:
        return ""
    header = [_clean_cell(c.text) for c in rows[0].cells]
    out = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        out.append("| " + " | ".join(_clean_cell(c.text) for c in row.cells) + " |")
    return "\n".join(out)


def _clean_cell(text: str) -> str:
    """Make a cell value safe to drop inside a one-line markdown table cell."""
    return (
        (text or "")
        .replace("\r", "")
        .replace("\n", "<br>")
        .replace("|", "\\|")
        .strip()
    )


# ── charts ──────────────────────────────────────────────────────────────────


def _chart_md(chart: Any) -> str:
    """Chart underlying data → a markdown table. Only a chart whose category /
    series *structure* can't be read degrades to ``[unsupported chart]``; a
    ragged series (fewer OR more values than categories) yields blank cells /
    blank categories rather than dropping data or discarding the whole chart."""
    try:
        plot = chart.plots[0]
        categories = [str(c) for c in plot.categories]
        # Pre-read each series' values defensively so one bad series can't sink
        # the whole table.
        series_data: list[tuple[str, list[Any]]] = []
        for s in chart.series:
            try:
                vals = list(s.values)
            except (ValueError, TypeError):
                vals = []
            series_data.append((str(s.name), vals))
    except (ValueError, IndexError, KeyError, AttributeError) as e:
        log.debug("pptx: unsupported chart (%s)", e)
        return "[unsupported chart]"

    header = ["Category"] + [_clean_cell(name) for name, _ in series_data]
    out = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    # Emit a row per category AND per overflow value — a series longer than the
    # category axis must not silently lose data (overflow rows get a blank
    # category label).
    n_rows = max([len(categories)] + [len(vals) for _, vals in series_data], default=0)
    for i in range(n_rows):
        cat = _clean_cell(categories[i]) if i < len(categories) else ""
        row = [cat]
        for _, vals in series_data:
            row.append(_clean_cell(_fmt_num(vals[i])) if i < len(vals) else "")
        out.append("| " + " | ".join(row) + " |")
    table_md = "\n".join(out)

    title = ""
    if chart.has_title:
        try:
            title = chart.chart_title.text_frame.text.strip()
        except (AttributeError, ValueError):
            title = ""
    heading = f"### Chart: {title}" if title else "### Chart"
    return f"{heading}\n\n{table_md}"


def _fmt_num(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


# ── pictures ──────────────────────────────────────────────────────────────────


def _is_picture(shape: Any) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    # A placeholder can carry an image; accessing .image raises if it doesn't.
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        try:
            return shape.image is not None
        except (KeyError, ValueError, AttributeError):
            return False
    return False


def _picture_md(shape: Any, opts: _Opts) -> str:
    desc = ""
    if opts.describe_images:
        blob = _picture_blob(shape)
        if blob is not None:
            try:
                desc = _describe_image(
                    blob, sensitivity=opts.sensitivity, ollama_client=opts.ollama_client
                )
            except PptxConversionError:
                # The sensitivity invariant (non-local lane) — never fall back,
                # never leak: propagate and abort the conversion.
                raise
            except Exception as e:  # noqa: BLE001 — a caption call failure
                # (Ollama down, model error, timeout) degrades to alt text rather
                # than sinking the whole deck.
                log.debug("pptx: image caption failed, falling back to alt (%s)", e)
                desc = ""
    if not desc:
        desc = _shape_alt_text(shape) or "image"
    return f"![{_escape_inline(desc)}]"


def _picture_blob(shape: Any) -> bytes | None:
    try:
        return shape.image.blob
    except (KeyError, ValueError, AttributeError):
        return None


def _shape_alt_text(shape: Any) -> str:
    """The ``descr`` (alt-text) attribute on the shape's ``cNvPr`` element."""
    try:
        cNvPr = shape._element.find(".//" + qn("p:cNvPr"))
        if cNvPr is not None:
            return (cNvPr.get("descr") or "").strip()
    except Exception:  # noqa: BLE001 — alt text is best-effort, never fatal
        return ""
    return ""


# ── image description (sensitivity-aware routing) ─────────────────────────────


_LOOPBACK_HOSTNAMES = {"localhost"}


def _is_loopback_url(base_url: str) -> bool:
    """True if ``base_url``'s host is loopback (``localhost`` or a 127.0.0.0/8 /
    ``::1`` address) — i.e. the request stays on this machine."""
    host = urlparse(base_url).hostname
    if host is None:
        return False
    if host in _LOOPBACK_HOSTNAMES:
        return True
    try:
        return ipaddress.ip_address(host).is_loopback
    except ValueError:
        return False


def _proxy_would_intercept(base_url: str) -> bool:
    """True if the environment's proxy settings would route a request to
    ``base_url`` through a proxy. ``OllamaClient`` calls ``requests`` with its
    default proxy-aware transport, so a configured HTTP(S)_PROXY/ALL_PROXY would
    carry a loopback request's body to a remote host unless NO_PROXY exempts it.
    We resolve this with requests' own logic so the answer matches what the
    client will actually do, and fail CLOSED (treat as intercepted) if we can't
    tell — image bytes must never risk leaving the box (CLAUDE.md
    #no-mnpi-to-cloud — was cited as §5.4)."""
    try:
        import requests.utils  # local import: requests is only needed on this path

        return bool(requests.utils.get_environ_proxies(base_url, no_proxy=None))
    except Exception:  # noqa: BLE001 — can't prove it's safe → assume it isn't
        return True


def _describe_image(blob: bytes, *, sensitivity: Sensitivity, ollama_client: OllamaClient) -> str:
    """Caption an image via ANTON's sensitivity-aware router.

    ``multimodal-extraction`` always resolves to the LOCAL ``gemma4:e4b`` lane,
    so a deck's image bytes never leave the box — for any sensitivity. We still
    assert the resolved lane is local as defence-in-depth: if routing policy
    ever changed, we refuse rather than silently leak an image to cloud
    (CLAUDE.md #no-mnpi-to-cloud, was cited as §5.4 — absolute). The client is
    a concrete ``OllamaClient`` whose endpoint the caller has already pinned to
    loopback, so the destination is the local Ollama lane."""
    lane = pick_lane("multimodal-extraction", sensitivity)
    provider, model = lane_to_model(lane)
    if provider != "ollama":
        raise PptxConversionError(
            f"sensitivity={sensitivity!r} resolved to non-local lane {lane!r}; "
            "refusing to send a slide image off-box"
        )
    b64 = base64.b64encode(blob).decode("ascii")
    resp = ollama_client.chat(
        model=model, prompt=_IMAGE_DESCRIBE_PROMPT, images=[b64], temperature=0.1
    )
    return (getattr(resp, "content", "") or "").strip()
